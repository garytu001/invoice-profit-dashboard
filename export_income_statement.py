"""
損益表匯出：Excel 和 PPT
參考台灣上市公司財報格式
"""
import io
from datetime import datetime


def export_income_statement_excel(IS: dict) -> bytes:
    """
    匯出損益表為 Excel
    格式參考台灣財報慣例
    """
    import openpyxl
    from openpyxl.styles import (
        Font, Alignment, PatternFill, Border, Side, numbers
    )
    
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "綜合損益表"
    
    # 欄寬設定
    ws.column_dimensions["A"].width = 35
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 12
    
    # 樣式定義
    header_font = Font(name="微軟正黑體", bold=True, size=14, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor="0F766E")
    title_font = Font(name="微軟正黑體", bold=True, size=11)
    bold_font = Font(name="微軟正黑體", bold=True, size=10)
    normal_font = Font(name="微軟正黑體", size=10)
    section_fill = PatternFill("solid", fgColor="E8F4F1")
    subtotal_fill = PatternFill("solid", fgColor="F0F9F8")
    center = Alignment(horizontal="center")
    right = Alignment(horizontal="right")
    thin = Side(style="thin", color="CCCCCC")
    border = Border(bottom=Side(style="thin", color="CCCCCC"))
    
    def fmt(val):
        if val is None:
            return "—"
        if val < 0:
            return f"({abs(val):,.0f})"
        return f"{val:,.0f}"
    
    def pct(val):
        if val is None:
            return "—"
        return f"{val*100:.1f}%"
    
    row = 1
    
    # 標題列
    ws.merge_cells(f"A{row}:C{row}")
    ws[f"A{row}"] = "綜合損益表"
    ws[f"A{row}"].font = header_font
    ws[f"A{row}"].fill = header_fill
    ws[f"A{row}"].alignment = center
    ws.row_dimensions[row].height = 30
    row += 1
    
    # 副標題
    ws.merge_cells(f"A{row}:C{row}")
    period_label = IS.get("year_month") or IS.get("period") or ""
    ws[f"A{row}"] = f"報表期間：{period_label}　　產生日期：{datetime.now().strftime('%Y-%m-%d')}"
    ws[f"A{row}"].font = normal_font
    ws[f"A{row}"].alignment = center
    row += 1
    
    # 空白列
    row += 1
    
    # 欄位標題
    headers = ["項目", "金額（元）", "佔營收比"]
    for col, h in enumerate(headers, 1):
        ws.cell(row=row, column=col, value=h).font = title_font
        ws.cell(row=row, column=col).fill = PatternFill("solid", fgColor="2C7873")
        ws.cell(row=row, column=col).font = Font(
            name="微軟正黑體", bold=True, size=10, color="FFFFFF"
        )
        ws.cell(row=row, column=col).alignment = center
    row += 1
    
    # 損益表資料列
    revenue = IS["revenue"]
    
    def write_row(label, val, pct_val, is_bold=False, fill=None):
        nonlocal row
        ws.cell(row=row, column=1, value=label)
        ws.cell(row=row, column=2, value=fmt(val))
        ws.cell(row=row, column=3, value=pct_val)
        for col in range(1, 4):
            cell = ws.cell(row=row, column=col)
            cell.font = bold_font if is_bold else normal_font
            cell.alignment = right if col > 1 else Alignment()
            if fill:
                cell.fill = fill
            cell.border = border
        row += 1
    
    write_row("營業收入淨額", revenue, pct(1.0), is_bold=True, fill=section_fill)
    write_row("　減：銷貨成本", -IS["cogs"], pct(-IS["cogs"]/revenue if revenue else 0))
    write_row("毛利（毛損）", IS["gross_profit"], pct(IS["gross_margin_rate"]), is_bold=True, fill=subtotal_fill)
    
    row += 1  # 空白
    write_row("營業費用", None, None, is_bold=True, fill=section_fill)
    
    for cat, amt in IS["opex_by_category"].items():
        write_row(f"　{cat}", -amt, pct(-amt/revenue if revenue else 0))
    
    write_row("　營業費用合計", -IS["total_opex"], pct(-IS["total_opex"]/revenue if revenue else 0), is_bold=True)
    write_row("營業利益（損失）", IS["operating_income"], pct(IS["operating_margin_rate"]), is_bold=True, fill=subtotal_fill)
    
    row += 1
    write_row("營業外收入及支出", None, None, is_bold=True, fill=section_fill)
    write_row("　其他收入", IS["non_operating_income"], "—")
    write_row("　其他支出", -IS["non_operating_expense"], "—")
    write_row("稅前淨利（損失）", IS["pretax_income"], pct(IS["net_margin_rate"]), is_bold=True, fill=subtotal_fill)
    write_row("　所得稅費用", -IS["income_tax"], "—")
    
    # 最終淨利（特別強調）
    ws.cell(row=row, column=1, value="本期淨利（損失）")
    ws.cell(row=row, column=2, value=fmt(IS["net_income"]))
    ws.cell(row=row, column=3, value=pct(IS["net_margin_rate"]))
    for col in range(1, 4):
        cell = ws.cell(row=row, column=col)
        cell.font = Font(name="微軟正黑體", bold=True, size=11, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="0F766E")
        cell.alignment = right if col > 1 else Alignment()
    row += 2
    
    # 附註區
    ws.cell(row=row, column=1, value="附註揭露").font = title_font
    row += 1
    
    ws.cell(row=row, column=1, value="附註一：重要會計政策").font = bold_font
    row += 1
    ws.merge_cells(f"A{row}:C{row}")
    ws.cell(row=row, column=1, value="本報表採用標準成本法估算銷貨成本，以請款單金額作為營業收入基礎。").font = normal_font
    row += 2
    
    ws.cell(row=row, column=1, value="附註二：客戶營收集中度").font = bold_font
    row += 1
    if IS.get("by_customer"):
        for c in sorted(IS["by_customer"], key=lambda x: x["revenue"], reverse=True)[:5]:
            pct_rev = c["revenue"] / revenue * 100 if revenue else 0
            ws.cell(row=row, column=1, value=f"  {c['customer_name']}")
            ws.cell(row=row, column=2, value=fmt(c["revenue"]))
            ws.cell(row=row, column=3, value=f"{pct_rev:.1f}%")
            for col in range(1, 4):
                ws.cell(row=row, column=col).font = normal_font
            row += 1
    
    # 儲存到 bytes
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def export_income_statement_ppt(IS: dict) -> bytes:
    """
    匯出損益表為 PPT
    簡報風格，適合向管理層報告
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 顏色定義
    BRAND = RGBColor(0x0F, 0x76, 0x6E)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0xF5, 0xF5, 0xF5)
    RED = RGBColor(0xC0, 0x39, 0x2B)
    GOLD = RGBColor(0xD4, 0xA0, 0x17)
    
    def add_textbox(slide, text, left, top, width, height,
                    font_size=12, bold=False, color=DARK,
                    align=PP_ALIGN.LEFT, bg_color=None):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if bg_color:
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox
    
    def fmt(val):
        if val is None:
            return "—"
        if val < 0:
            return f"({abs(val):,.0f})"
        return f"{val:,.0f}"
    
    def pct(val):
        if val is None:
            return "—"
        return f"{val*100:.1f}%"
    
    revenue = IS["revenue"]
    
    # === 第一張：封面 ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    
    # 背景色
    bg = slide1.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    
    add_textbox(slide1, "綜合損益表", 1, 2, 11, 1.5,
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide1, "Comprehensive Income Statement",
                1, 3.3, 11, 0.8,
                font_size=16, color=RGBColor(0xA0, 0xA0, 0xB0),
                align=PP_ALIGN.CENTER)
    period_label = IS.get("year_month") or IS.get("period") or ""
    add_textbox(slide1, f"報表期間：{period_label}",
                1, 4.5, 11, 0.6,
                font_size=14, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide1, f"產生日期：{datetime.now().strftime('%Y-%m-%d')}",
                1, 5.1, 11, 0.5,
                font_size=12, color=RGBColor(0x80, 0x80, 0x90),
                align=PP_ALIGN.CENTER)
    
    # === 第二張：KPI 摘要 ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_textbox(slide2, "財務摘要", 0.5, 0.3, 12, 0.8,
                font_size=24, bold=True, color=BRAND)
    
    # 四個 KPI 卡片
    kpis = [
        ("營業收入", f"${revenue:,.0f}", BRAND),
        ("毛利率", pct(IS["gross_margin_rate"]),
         BRAND if IS["gross_margin_rate"] and IS["gross_margin_rate"] > 0 else RED),
        ("營業利益率", pct(IS["operating_margin_rate"]),
         BRAND if IS["operating_income"] > 0 else RED),
        ("本期淨利率", pct(IS["net_margin_rate"]),
         BRAND if IS["net_income"] > 0 else RED),
    ]
    
    for i, (label, val, color) in enumerate(kpis):
        x = 0.5 + i * 3.2
        # 卡片背景
        shape = slide2.shapes.add_shape(
            1, Inches(x), Inches(1.3), Inches(2.9), Inches(1.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = GRAY
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        
        add_textbox(slide2, label, x+0.1, 1.4, 2.7, 0.5,
                    font_size=11, color=RGBColor(0x60, 0x60, 0x70))
        add_textbox(slide2, val, x+0.1, 1.9, 2.7, 0.9,
                    font_size=22, bold=True, color=color)
    
    # 損益結構長條（簡化版）
    add_textbox(slide2, "損益結構", 0.5, 3.3, 12, 0.5,
                font_size=14, bold=True, color=DARK)
    
    items_bar = [
        ("營業收入", revenue, BRAND),
        ("銷貨成本", IS["cogs"], RGBColor(0xEF, 0x88, 0x44)),
        ("毛利", IS["gross_profit"], RGBColor(0x22, 0xC5, 0x5E)),
        ("營業費用", IS["total_opex"], RGBColor(0xEF, 0x44, 0x44)),
        ("營業利益", IS["operating_income"],
         RGBColor(0x22, 0xC5, 0x5E) if IS["operating_income"] >= 0 else RED),
        ("本期淨利", IS["net_income"],
         RGBColor(0x22, 0xC5, 0x5E) if IS["net_income"] >= 0 else RED),
    ]
    
    max_val = max(abs(v) for _, v, _ in items_bar) or 1
    bar_max_w = 7.0
    
    for i, (label, val, color) in enumerate(items_bar):
        y = 3.9 + i * 0.52
        add_textbox(slide2, label, 0.5, y, 2.2, 0.45, font_size=10, color=DARK)
        bar_w = abs(val) / max_val * bar_max_w
        shape = slide2.shapes.add_shape(
            1, Inches(2.8), Inches(y+0.05), Inches(max(bar_w, 0.05)), Inches(0.32)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        add_textbox(slide2, fmt(val), 2.8 + bar_w + 0.1, y, 2.5, 0.45,
                    font_size=10, bold=True, color=DARK)
    
    # === 第三張：完整損益表 ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_textbox(slide3, "綜合損益表明細", 0.5, 0.2, 12, 0.6,
                font_size=20, bold=True, color=BRAND)
    
    # 表格資料
    table_data = [
        ("項目", "金額（元）", "佔營收比", True, BRAND, WHITE),
        ("營業收入淨額", fmt(revenue), pct(1.0), True, GRAY, DARK),
        ("　減：銷貨成本", fmt(-IS["cogs"]),
         pct(-IS["cogs"]/revenue if revenue else 0), False, WHITE, DARK),
        ("毛利（毛損）", fmt(IS["gross_profit"]),
         pct(IS["gross_margin_rate"]), True,
         RGBColor(0xE8, 0xF4, 0xF1), DARK),
        ("營業費用合計", fmt(-IS["total_opex"]),
         pct(-IS["total_opex"]/revenue if revenue else 0), False, WHITE, DARK),
        ("營業利益（損失）", fmt(IS["operating_income"]),
         pct(IS["operating_margin_rate"]), True,
         RGBColor(0xE8, 0xF4, 0xF1), DARK),
        ("稅前淨利（損失）", fmt(IS["pretax_income"]),
         pct(IS["net_margin_rate"]), True,
         RGBColor(0xE8, 0xF4, 0xF1), DARK),
        ("本期淨利（損失）", fmt(IS["net_income"]),
         pct(IS["net_margin_rate"]), True, BRAND, WHITE),
    ]
    
    for i, (label, val, pct_v, bold, bg, fg) in enumerate(table_data):
        y = 1.0 + i * 0.62
        for col_idx, (text, x, w) in enumerate([
            (label, 0.5, 6.5),
            (val, 7.2, 2.8),
            (pct_v, 10.2, 2.5),
        ]):
            shape = slide3.shapes.add_shape(
                1, Inches(x), Inches(y), Inches(w), Inches(0.55)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            add_textbox(slide3, text, x+0.05, y+0.07, w-0.1, 0.42,
                        font_size=10 if not bold else 11,
                        bold=bold, color=fg,
                        align=PP_ALIGN.RIGHT if col_idx > 0 else PP_ALIGN.LEFT)
    
    # === 第四張：附註 ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, "附註揭露", 0.5, 0.2, 12, 0.6,
                font_size=20, bold=True, color=BRAND)
    
    add_textbox(slide4, "附註一：重要會計政策", 0.5, 1.0, 12, 0.5,
                font_size=13, bold=True, color=DARK)
    add_textbox(slide4,
                "本報表採用標準成本法估算銷貨成本，以請款單金額作為營業收入基礎。\n"
                "單位不一致或缺少成本資料之明細不計入成本計算。",
                0.5, 1.5, 12, 0.8, font_size=11, color=DARK)
    
    add_textbox(slide4, "附註二：客戶營收集中度（前五大）",
                0.5, 2.5, 12, 0.5, font_size=13, bold=True, color=DARK)
    
    if IS.get("by_customer"):
        top5 = sorted(IS["by_customer"], key=lambda x: x["revenue"], reverse=True)[:5]
        for i, c in enumerate(top5):
            pct_rev = c["revenue"] / revenue * 100 if revenue else 0
            add_textbox(
                slide4,
                f"{i+1}. {c['customer_name']}　${c['revenue']:,.0f}　"
                f"（{pct_rev:.1f}%）",
                0.8, 3.0 + i * 0.45, 12, 0.42,
                font_size=11, color=DARK
            )
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()